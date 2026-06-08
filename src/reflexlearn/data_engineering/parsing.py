"""多格式文档解析：pdf / docx / pptx / html / md / txt → 结构化 sections（缺库即降级）。

设计：
- 二进制格式（pdf/docx/pptx）与 html 的解析库在函数内 import，未安装时抛 ParserUnavailable，
  由 ingest 上层捕获 → 该文档 degraded（chunks=0）但写链路仍 200（降级铁律）。
- md/txt 用内置解码，纯文本永不降级。
- 结构感知：尽量切出 (heading, text) 段，供 chunking 在 section 边界内聚合，保留文档结构语义。
"""
from __future__ import annotations

from dataclasses import dataclass, field


class ParserUnavailable(RuntimeError):
    """解析依赖缺失 / 文件打开失败 / 不支持的格式——上层据此对该文档降级。"""


@dataclass
class Section:
    heading: str
    text: str


@dataclass
class ParsedDoc:
    title: str
    format: str
    sections: list[Section] = field(default_factory=list)

    @property
    def text(self) -> str:
        return "\n\n".join(s.text for s in self.sections if s.text)


def _basename(filename: str) -> str:
    return (filename or "").replace("\\", "/").split("/")[-1]


def _stem(filename: str) -> str:
    name = _basename(filename)
    return name.rsplit(".", 1)[0] if "." in name else name


def _ext(filename: str) -> str:
    name = _basename(filename).lower()
    return "." + name.rsplit(".", 1)[-1] if "." in name else ""


def _first_line(text: str) -> str:
    for ln in (text or "").split("\n"):
        s = ln.strip().lstrip("#").strip()
        if s:
            return s
    return ""


def detect_format(filename: str) -> str:
    """返回不带点的格式名（解析失败时 ingest 用于 IngestResult.format）。"""
    ext = _ext(filename).lstrip(".")
    return {"markdown": "md", "htm": "html"}.get(ext, ext)


def parse_document(filename: str, raw: bytes) -> ParsedDoc:
    """按扩展名分派解析。不支持的扩展名抛 ParserUnavailable。"""
    ext = _ext(filename)
    if ext in (".md", ".markdown"):
        return _parse_markdown(filename, raw)
    if ext == ".txt":
        return _parse_text(filename, raw)
    if ext in (".html", ".htm"):
        return _parse_html(filename, raw)
    if ext == ".pdf":
        return _parse_pdf(filename, raw)
    if ext == ".docx":
        return _parse_docx(filename, raw)
    if ext == ".pptx":
        return _parse_pptx(filename, raw)
    raise ParserUnavailable(f"unsupported format: {ext or '(none)'}")


# ------------------------------------------------------------------ 纯文本 / markdown


def _parse_markdown(filename: str, raw: bytes) -> ParsedDoc:
    text = raw.decode("utf-8", errors="replace")
    sections: list[Section] = []
    title = ""
    cur_heading = ""
    cur_lines: list[str] = []
    for line in text.split("\n"):
        if line.lstrip().startswith("#"):
            body = "\n".join(cur_lines).strip()
            if body:
                sections.append(Section(cur_heading, body))
            cur_heading = line.lstrip().lstrip("#").strip()
            cur_lines = []
            if not title:
                title = cur_heading
        else:
            cur_lines.append(line)
    body = "\n".join(cur_lines).strip()
    if body:
        sections.append(Section(cur_heading, body))
    if not sections and text.strip():
        sections = [Section("", text.strip())]
    return ParsedDoc(title=title or _first_line(text) or _stem(filename), format="md", sections=sections)


def _parse_text(filename: str, raw: bytes) -> ParsedDoc:
    text = raw.decode("utf-8", errors="replace").strip()
    sections = [Section("", text)] if text else []
    return ParsedDoc(title=_first_line(text) or _stem(filename), format="txt", sections=sections)


# ------------------------------------------------------------------ html


def _parse_html(filename: str, raw: bytes) -> ParsedDoc:
    try:
        from bs4 import BeautifulSoup
    except Exception as e:  # 缺 beautifulsoup4
        raise ParserUnavailable(f"beautifulsoup4 not available: {e}") from e
    soup = BeautifulSoup(raw, "html.parser")
    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    if not title:
        h1 = soup.find("h1")
        if h1:
            title = h1.get_text(strip=True)
    text = soup.get_text("\n", strip=True)
    sections = [Section("", text)] if text else []
    return ParsedDoc(title=title or _stem(filename), format="html", sections=sections)


# ------------------------------------------------------------------ pdf / docx / pptx


def _parse_pdf(filename: str, raw: bytes) -> ParsedDoc:
    try:
        import fitz  # PyMuPDF
    except Exception as e:
        raise ParserUnavailable(f"pymupdf not available: {e}") from e
    try:
        doc = fitz.open(stream=raw, filetype="pdf")
    except Exception as e:
        raise ParserUnavailable(f"pdf open failed: {e}") from e
    title = ""
    try:
        title = (doc.metadata or {}).get("title", "") or ""
    except Exception:
        pass
    sections: list[Section] = []
    for i, page in enumerate(doc):
        try:
            t = page.get_text().strip()
        except Exception:
            t = ""
        if t:
            sections.append(Section(f"第 {i + 1} 页", t))
    doc.close()
    return ParsedDoc(title=title.strip() or _stem(filename), format="pdf", sections=sections)


def _parse_docx(filename: str, raw: bytes) -> ParsedDoc:
    try:
        import docx  # python-docx
    except Exception as e:
        raise ParserUnavailable(f"python-docx not available: {e}") from e
    import io

    try:
        d = docx.Document(io.BytesIO(raw))
    except Exception as e:
        raise ParserUnavailable(f"docx open failed: {e}") from e
    title = ""
    try:
        title = d.core_properties.title or ""
    except Exception:
        pass
    sections: list[Section] = []
    cur_heading = ""
    cur_lines: list[str] = []
    for p in d.paragraphs:
        ptext = (p.text or "").strip()
        style = getattr(p.style, "name", "") or ""
        if style.startswith("Heading") or style == "Title":
            if cur_lines:
                sections.append(Section(cur_heading, "\n".join(cur_lines)))
                cur_lines = []
            cur_heading = ptext
            if not title and ptext:
                title = ptext
        elif ptext:
            cur_lines.append(ptext)
    if cur_lines:
        sections.append(Section(cur_heading, "\n".join(cur_lines)))
    if not sections:
        full = "\n".join((p.text or "") for p in d.paragraphs).strip()
        sections = [Section("", full)] if full else []
    return ParsedDoc(title=(title or _stem(filename)).strip(), format="docx", sections=sections)


def _parse_pptx(filename: str, raw: bytes) -> ParsedDoc:
    try:
        from pptx import Presentation  # python-pptx
    except Exception as e:
        raise ParserUnavailable(f"python-pptx not available: {e}") from e
    import io

    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as e:
        raise ParserUnavailable(f"pptx open failed: {e}") from e
    sections: list[Section] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        body = "\n".join(parts)
        if body:
            sections.append(Section(f"幻灯片 {i + 1}", body))
    return ParsedDoc(title=_stem(filename), format="pptx", sections=sections)
