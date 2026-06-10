from __future__ import annotations

import builtins
import io

import pytest

from reflexlearn.data_engineering.parsing import (
    ParserUnavailable,
    detect_format,
    parse_document,
)


def test_parse_markdown_sections_and_title():
    raw = "# 机器学习\n\n## 线性回归\n\n最小二乘法求参数\n\n## 梯度下降\n\n迭代最小化损失".encode()
    doc = parse_document("ml.md", raw)
    assert doc.format == "md"
    assert doc.title == "机器学习"
    headings = [s.heading for s in doc.sections]
    assert "线性回归" in headings and "梯度下降" in headings
    assert "最小二乘法求参数" in doc.text


def test_parse_text_single_section():
    doc = parse_document("note.txt", "第一行\n第二行".encode())
    assert doc.format == "txt"
    assert len(doc.sections) == 1
    assert "第二行" in doc.sections[0].text


def test_parse_markdown_no_heading_fallback():
    doc = parse_document("plain.md", "没有标题的正文内容".encode())
    assert doc.sections
    assert "没有标题的正文内容" in doc.text


def test_detect_format():
    assert detect_format("a.MD") == "md"
    assert detect_format("a.markdown") == "md"
    assert detect_format("a.htm") == "html"
    assert detect_format("a.pdf") == "pdf"


def test_unsupported_format_raises():
    with pytest.raises(ParserUnavailable):
        parse_document("a.xyz", b"data")


def test_parse_html():
    pytest.importorskip("bs4")
    raw = (
        "<html><head><title>标题T</title></head>"
        "<body><h1>大标题</h1><p>HTML正文</p><script>var x=1;</script></body></html>"
    ).encode()
    doc = parse_document("p.html", raw)
    assert doc.format == "html"
    assert doc.title == "标题T"
    assert "HTML正文" in doc.text
    assert "var x" not in doc.text  # script 已去除


def test_parse_pdf():
    fitz = pytest.importorskip("fitz")
    d = fitz.open()
    page = d.new_page()
    page.insert_text((72, 72), "linear regression notes")
    raw = d.tobytes()
    d.close()
    doc = parse_document("p.pdf", raw)
    assert doc.format == "pdf"
    assert "linear regression" in doc.text


def test_parse_docx():
    docx = pytest.importorskip("docx")
    d = docx.Document()
    d.add_heading("文档标题", level=1)
    d.add_paragraph("文档正文内容")
    b = io.BytesIO()
    d.save(b)
    doc = parse_document("p.docx", b.getvalue())
    assert doc.format == "docx"
    assert "文档正文内容" in doc.text


def test_parse_pptx():
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "幻灯片文本内容"
    b = io.BytesIO()
    prs.save(b)
    doc = parse_document("p.pptx", b.getvalue())
    assert doc.format == "pptx"
    assert "幻灯片文本内容" in doc.text


def test_pdf_missing_lib_degrades(monkeypatch):
    """缺 pymupdf → ParserUnavailable（上层据此降级）。"""
    real_import = builtins.__import__

    def fake_import(name, *a, **k):
        if name == "fitz":
            raise ImportError("no fitz")
        return real_import(name, *a, **k)

    monkeypatch.setattr(builtins, "__import__", fake_import)
    with pytest.raises(ParserUnavailable):
        parse_document("p.pdf", b"%PDF-1.4 fake")
